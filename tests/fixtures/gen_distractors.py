"""Generate 4 off-topic distractor docs in 4 formats, each a DIFFERENT subject,
all unrelated to the 'AI Customer Agent' knowledge base. Ingesting these lets eval
check retrieval doesn't pull cross-topic chunks into an on-topic answer.

Run: python -m tests.fixtures.gen_distractors   (writes files next to this script)
"""
from pathlib import Path

OUT = Path(__file__).resolve().parent


def docx_coffee(path: Path) -> None:
    from docx import Document
    d = Document()
    d.add_heading("Home Coffee Brewing Guide", 0)
    d.add_paragraph(
        "Great coffee starts with fresh beans and the right water temperature. "
        "Water between 90 and 96 degrees Celsius extracts flavor without scorching."
    )
    d.add_heading("Pour-Over Method", level=1)
    d.add_paragraph(
        "Use a 1:16 coffee-to-water ratio. For 20 grams of medium-ground coffee, "
        "pour 320 grams of water in slow circles over about three minutes. "
        "Bloom the grounds with 40 grams of water for 30 seconds first."
    )
    d.add_heading("French Press", level=1)
    d.add_paragraph(
        "Steep coarse grounds for four minutes, then press slowly. "
        "A coarser grind avoids a muddy, over-extracted cup."
    )
    d.save(str(path))


def pdf_hiking(path: Path) -> None:
    from fpdf import FPDF
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 16)
    pdf.cell(0, 10, "Cascade Ridge Trail Guide", new_x="LMARGIN", new_y="NEXT")
    pdf.set_font("Helvetica", size=11)
    body = (
        "The Cascade Ridge Trail is a 12.4 km loop with 850 meters of elevation gain. "
        "The trailhead sits at Pine Hollow parking lot, open from dawn to dusk.\n\n"
        "Difficulty: Moderate to hard. Allow five to six hours round trip.\n"
        "Water: Two natural springs at km 3 and km 7; treat before drinking.\n"
        "Best season: Late June through September, before the first snow.\n\n"
        "Wildlife: Black bears are active in spring. Carry your food in a bear "
        "canister and make noise on blind corners."
    )
    pdf.multi_cell(0, 7, body)
    pdf.output(str(path))


def xlsx_garden(path: Path) -> None:
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Watering Schedule"
    ws.append(["Plant", "Location", "Water Every (days)", "Sunlight", "Notes"])
    rows = [
        ["Basil", "Kitchen windowsill", 2, "Full sun", "Pinch flowers to keep leaves growing"],
        ["Snake Plant", "Living room corner", 14, "Low light", "Drought tolerant, do not overwater"],
        ["Tomato", "Balcony", 1, "Full sun", "Needs support cage once 40 cm tall"],
        ["Fern", "Bathroom shelf", 3, "Indirect", "Likes humidity, mist between waterings"],
        ["Cactus", "South window", 21, "Full sun", "Water deeply but rarely"],
    ]
    for r in rows:
        ws.append(r)
    wb.save(str(path))


def pptx_rome(path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    title = prs.slides.add_slide(prs.slide_layouts[0])
    title.shapes.title.text = "The Roman Republic"
    title.placeholders[1].text = "From founding to empire, 509-27 BCE"

    bullet = prs.slides.add_slide(prs.slide_layouts[1])
    bullet.shapes.title.text = "Key Events"
    tf = bullet.placeholders[1].text_frame
    tf.text = "509 BCE: Republic founded after the last king is expelled"
    for line in [
        "264-146 BCE: Three Punic Wars against Carthage",
        "49 BCE: Julius Caesar crosses the Rubicon",
        "44 BCE: Caesar assassinated on the Ides of March",
        "27 BCE: Augustus becomes first emperor, ending the Republic",
    ]:
        tf.add_paragraph().text = line
    prs.save(str(path))


SPECS = [
    ("coffee_brewing_guide.docx", docx_coffee),
    ("cascade_ridge_trail.pdf", pdf_hiking),
    ("garden_watering_schedule.xlsx", xlsx_garden),
    ("roman_republic_history.pptx", pptx_rome),
]


def main() -> list[Path]:
    paths = []
    for name, fn in SPECS:
        p = OUT / name
        fn(p)
        assert p.exists() and p.stat().st_size > 0, f"failed to write {name}"
        paths.append(p)
        print(f"wrote {name} ({p.stat().st_size} bytes)")
    return paths


if __name__ == "__main__":
    main()
    print("OK")
